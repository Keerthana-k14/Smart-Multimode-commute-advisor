from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Define some useful layouts
    title_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]
    
    # Helper to add a slide with title and bullets
    def add_bullet_slide(title, points):
        slide = prs.slides.add_slide(bullet_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = title
        tf = body_shape.text_frame
        
        for i, point in enumerate(points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(18)

    # 1. Title Slide
    slide = prs.slides.add_slide(title_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "GlassEngine"
    subtitle.text = "A Self-Adapting System for Explaining Digital Authenticity\nShifting Detection from AI Pattern Recognition -> Mathematical Physics Validation"

    # 2. INTRODUCTION
    add_bullet_slide(
        "01 / INTRODUCTION: The Erosion of Digital Trust",
        [
            "The Paradigm Shift: Generative AI (GANs, Diffusion Models) has evolved from obvious visual errors to synthesizing hyper-realistic media that completely bypasses human perception.",
            "The Arms Race: Traditional deepfake detection fought AI with AI, creating a costly, endless cycle where generators continuously learn to evade neural detectors.",
            "Our Domain: Operating at the intersection of Digital Forensics, Image Integrity Analysis, and Generative AI Security—targeting the physical impossibilities AI cannot fake."
        ]
    )

    # 3. INTRODUCTION 2
    add_bullet_slide(
        "02 / INTRODUCTION: Introducing Physical Forensics",
        [
            "THE FLAW IN AI: AI algorithms synthesize images randomly out of thin air, failing to simulate invisible rules of the physical world.",
            "This leaves detectable artifacts: Missing Bayer sensor noise, perfect entropy, phase mismatches, no hardware fingerprint.",
            "OUR SOLUTION: GlassEngine—a proprietary mathematical verification pipeline treating images as biological specimens.",
            "Digital DNA Extraction includes: Sensor noise analysis, Light phase extraction, Entropy measurement, Processing history forensics."
        ]
    )

    # 4. PROBLEM STATEMENT
    add_bullet_slide(
        "03 / PROBLEM STATEMENT",
        [
            "Current systems rely on continuously retraining neural networks—a reactive strategy that remains unexplainable, computationally exhaustive, and highly vulnerable.",
            "There is a critical need for an explainable, deterministic framework capable of validating image authenticity without exclusively relying on deep learning.",
            "Reactive: Neural nets retrain after every new generator release.",
            "Unexplainable: CNN 'black-box' guesses with zero forensic evidence.",
            "Fragile: Near 50% accuracy when generators change texture architecture."
        ]
    )

    # 5. OBJECTIVE
    add_bullet_slide(
        "04 / OBJECTIVE: Project Objectives",
        [
            "01 Multi-Pillar Framework: Construct a custom forensic engine analyzing distinct deterministic markers (Phase variances, Chaos distribution, Noise Variance, Processing history).",
            "02 Formulate Novel Mathematics: Develop a proprietary 'High-Frequency Enhanced Phase Extractor' and custom 'Collision Penalty Formula' to hunt synthetic micro-textures.",
            "03 Ensure Explainability: Design a real-time, side-by-side comparative dashboard that transitions deepfake detection to irrefutable mathematical proof."
        ]
    )

    # 6. LITERATURE SURVEY
    add_bullet_slide(
        "05 / LITERATURE SURVEY: Foundational Approaches",
        [
            "Vision-Language Feature Extraction (Cappelletti et al., 2024): Deep semantic encoding captures synthetic anomalies without massive datasets.",
            "Inter-Channel Phase Analysis (Khan et al., ICCK, 2025): Real hardware captures perfectly synced phases; AI generates scattered phase mismatches.",
            "Latent Information Density (f-InfoED, 2024): Generative AI is 'too efficient,' creating unnatural smoothness lacking nature's micro-chaos.",
            "Hybrid Forensic Noise Tracking (IJSAT, 2025): Synthetic imagery lacks continuous ambient sensor noise across dark areas."
        ]
    )

    # 7. LITERATURE SURVEY - GAPS
    add_bullet_slide(
        "06 / LITERATURE SURVEY — GAPS",
        [
            "GAP 01 The 'Black Box' Trap: Nearly 90% of detectors use giant CNNs that guess blindly, lacking explainable visual proof.",
            "GAP 02 Generative Escapism: Neural detectors drop to near 50% accuracy against zero-day models when AI generators change texture architecture.",
            "GAP 03 The High-Frequency Blindspot: Many models test global image metadata but fail to examine ultra-detailed structural boundaries where AI struggles."
        ]
    )

    # 8. METHODOLOGY
    add_bullet_slide(
        "07 / METHODOLOGY: GlassEngine Core Architecture",
        [
            "1. Input / UI Routing: Target image ingested through Single Neural or Dual Math Comparative routing nodes.",
            "2. Preprocessing Layer: Image structurally center-cropped to preserve critical hardware configurations.",
            "3. Diagnostic Filtering: Simultaneous routing through unsharp masking layers and custom entropy calculators.",
            "4. Verdict Dashboard: Phase variances mapped against sterile noise matrices—rendered dynamically to the user."
        ]
    )

    # 9. METHODOLOGY 2
    add_bullet_slide(
        "08 / METHODOLOGY: Proprietary Forensic Modules",
        [
            "M1 Custom HF-ICPC Engine: Applies architectural algorithm to emphasize structural boundaries and hunts for missing physical Bayer-locks.",
            "M2 GlassEngine Collision Penalty: Custom mathematical entropy check that instantly flags and penalizes heavily 'smooth' zones synthesized by deepfakes.",
            "M3 Sensor PRNU Extraction: High-pass filter execution targeting baseline silicon camera noise, unveiling mathematically 'sterile' shadow areas.",
            "M4 Zero-Shot Semantic Engine: Customized vision-language inference pipeline utilizing foundational Transformer backbones to capture edge-case perceptual hallucinations."
        ]
    )

    # 10. REFERENCES
    add_bullet_slide(
        "09 / REFERENCES",
        [
            "1. Cappelletti et al., 'CLIP & Few-Shot Learning for Forensics,' ICPR, 2024.",
            "2. Tan et al., 'Vision-Language Forensic Inference,' AAAI, 2025.",
            "3. Khan et al., 'Inter-Channel Phase Correlation in Bayer Filters,' ICCK, 2025.",
            "4. ArXiv Preprint, 'Latent Entropy in Generative Models (f-InfoED),' 2024.",
            "5. International Journal of Scientific Research, 'NoiseDF & Hybrid Forensic Frameworks,' 2025."
        ]
    )

    # Save presentation
    output_path = r"c:\Users\Keerthana K\OneDrive\Desktop\GlassEngine_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
